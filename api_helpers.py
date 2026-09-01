"""
api_helpers.py — shared integration helpers for all modules.

Consolidates the StaffMeetingBuilder baseline helpers (Outlook search,
Confluence pull, file parsing) and adds Jira (REST v2, Bearer PAT) and
additional MS Graph helpers (mail draft/send, Teams post, calendar).

HITL guardrail: send_mail() defaults to creating an Outlook *draft* and
post_teams_message() must be called explicitly from a confirm click — nothing
in this module is invoked automatically by the LLM.
"""
import io
import re
import time
import urllib3

import requests
import pandas as pd
import docx
import pptx
from bs4 import BeautifulSoup

from datetime import datetime, timedelta

import config
from ms_auth import get_ms_token

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

GRAPH_ROOT = "https://graph.microsoft.com/v1.0"


# --------------------------------------------------------------------------- #
# Generic helpers
# --------------------------------------------------------------------------- #
def clean_html(raw_html):
    if not raw_html:
        return ""
    soup = BeautifulSoup(raw_html, "html.parser")
    text = soup.get_text(separator="\n")
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def extract_text_from_file(uploaded_file):
    """Extract text from a Streamlit UploadedFile (txt/csv/xlsx/docx/pptx).
    PDFs are handled separately by callers (sent to Gemini as a Part)."""
    ext = uploaded_file.name.split(".")[-1].lower()
    try:
        if ext in ["txt", "csv"]:
            return uploaded_file.getvalue().decode("utf-8", errors="replace")
        elif ext == "xlsx":
            df = pd.read_excel(uploaded_file)
            return df.to_csv(index=False)
        elif ext == "docx":
            doc = docx.Document(uploaded_file)
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext == "pptx":
            prs = pptx.Presentation(uploaded_file)
            runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        runs.append(shape.text)
            return "\n".join(runs)
        return ""
    except Exception as e:
        return f"[Error extracting text from {uploaded_file.name}: {e}]"


def extract_text_from_bytes(filename, data):
    """Extract text from raw bytes (txt/csv/xlsx/docx/pptx/pdf), e.g. a downloaded
    Jira attachment. For PDFs, only the embedded text is pulled (images are
    ignored) so image-heavy reports don't blow up the token budget. Returns ""
    for unsupported types or image-only PDFs."""
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    try:
        if ext in ["txt", "csv"]:
            return data.decode("utf-8", errors="replace")
        elif ext == "xlsx":
            return pd.read_excel(io.BytesIO(data)).to_csv(index=False)
        elif ext == "docx":
            doc = docx.Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs]
            # Most structured content (e.g. goal-appraisal rows) lives in tables,
            # which are NOT in doc.paragraphs — pull each row as tab-joined cells.
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
            return "\n".join(p for p in parts if p and p.strip())
        elif ext == "pptx":
            prs = pptx.Presentation(io.BytesIO(data))
            runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        runs.append(shape.text)
            return "\n".join(runs)
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                return "[pypdf not installed — cannot extract PDF text]"
            reader = PdfReader(io.BytesIO(data))
            pages = [(page.extract_text() or "") for page in reader.pages]
            return "\n".join(p for p in pages if p.strip())
        return ""
    except Exception as e:
        return f"[Error extracting text from {filename}: {e}]"


# --------------------------------------------------------------------------- #
# MS Graph — Outlook mail
# --------------------------------------------------------------------------- #
def _graph_headers():
    return {
        "Authorization": f"Bearer {get_ms_token()}",
        "Content-Type": "application/json",
        "ConsistencyLevel": "eventual",
    }


def fetch_emails_from_graph(query_df, start_date, end_date):
    """Search Outlook for the configured senders/subjects/categories plus the
    'staff meeting include' category, within the date range. Returns a
    de-duplicated list of message dicts."""
    headers = _graph_headers()
    start_str = start_date.strftime("%Y-%m-%d")
    end_str = end_date.strftime("%Y-%m-%d")
    date_query = f"received:{start_str}..{end_str}"

    all_messages = []
    seen_ids = set()
    queries_to_run = [f'{date_query} AND category:"staff meeting include"']

    if query_df is not None:
        for _, row in query_df.iterrows():
            search_type = str(row.get("Search Type", "")).strip().lower()
            search_value = str(row.get("Value", "")).strip()
            if not search_value:
                continue
            if search_type == "sender":
                queries_to_run.append(f'{date_query} AND from:"{search_value}"')
            elif search_type == "subject":
                queries_to_run.append(f'{date_query} AND subject:"{search_value}"')
            elif search_type == "category":
                queries_to_run.append(f'{date_query} AND category:"{search_value}"')

    for kql_query in queries_to_run:
        safe_kql = kql_query.replace('"', '\\"')
        url = (
            f"{GRAPH_ROOT}/me/messages?"
            f'$search="{safe_kql}"&'
            "$select=id,subject,sender,body,receivedDateTime,categories"
        )
        while url:
            response = requests.get(url, headers=headers, verify=False)
            if response.status_code != 200:
                raise RuntimeError(
                    f"Graph API Error: {response.status_code} - {response.text}"
                )
            data = response.json()
            for msg in data.get("value", []):
                msg_id = msg.get("id")
                if msg_id not in seen_ids:
                    seen_ids.add(msg_id)
                    all_messages.append(msg)
            url = data.get("@odata.nextLink")
    return all_messages


def fetch_ooo_emails(start_date, end_date, keywords=("OOO", "VTO", "PTO", "out of office")):
    """Pull recent messages whose subject mentions OOO/VTO/PTO for the Module F
    availability digest. Returns a list of message dicts."""
    headers = _graph_headers()
    start_str = start_date.strftime("%Y-%m-%d")
    end_str = end_date.strftime("%Y-%m-%d")
    date_query = f"received:{start_str}..{end_str}"

    all_messages = []
    seen_ids = set()
    for kw in keywords:
        kql = f'{date_query} AND subject:"{kw}"'.replace('"', '\\"')
        url = (
            f"{GRAPH_ROOT}/me/messages?"
            f'$search="{kql}"&'
            "$select=id,subject,sender,body,receivedDateTime"
        )
        response = requests.get(url, headers=headers, verify=False)
        if response.status_code != 200:
            continue
        for msg in response.json().get("value", []):
            mid = msg.get("id")
            if mid not in seen_ids:
                seen_ids.add(mid)
                all_messages.append(msg)
    return all_messages


def fetch_ooo_requests(
    start_date, end_date, sender="", keywords=("OOO", "VTO", "PTO", "out of office")
):
    """Pull the HR system's OoO approval-request emails for the standalone OoO
    Management tool. When `sender` is set, search by that from-address (most
    reliable); otherwise fall back to the subject keywords. Returns de-duplicated
    dicts: {id, subject, sender_name, sender_addr, received, body(cleaned)}."""
    headers = _graph_headers()
    start_str = start_date.strftime("%Y-%m-%d")
    end_str = end_date.strftime("%Y-%m-%d")
    date_query = f"received:{start_str}..{end_str}"

    if sender.strip():
        queries = [f'{date_query} AND from:"{sender.strip()}"']
    else:
        queries = [f'{date_query} AND subject:"{kw}"' for kw in keywords]

    results = []
    seen_ids = set()
    for q in queries:
        kql = q.replace('"', '\\"')
        url = (
            f"{GRAPH_ROOT}/me/messages?"
            f'$search="{kql}"&'
            "$select=id,subject,sender,body,receivedDateTime"
        )
        response = requests.get(url, headers=headers, verify=False)
        if response.status_code != 200:
            continue
        for msg in response.json().get("value", []):
            mid = msg.get("id")
            if mid in seen_ids:
                continue
            seen_ids.add(mid)
            email = (msg.get("sender", {}) or {}).get("emailAddress", {}) or {}
            results.append(
                {
                    "id": mid,
                    "subject": msg.get("subject", "No Subject"),
                    "sender_name": email.get("name", "Unknown"),
                    "sender_addr": email.get("address", ""),
                    "received": msg.get("receivedDateTime", ""),
                    "body": clean_html(msg.get("body", {}).get("content", "")),
                }
            )
    return results


def fetch_recent_inbox(since_iso):
    """Fetch inbox messages received at/after `since_iso` (ISO-8601 UTC, e.g.
    '2026-06-15T00:00:00Z'), newest first, following pagination. Mirrors the
    EmailToJira inbox pull used by the Phase 0 Email Action Identifier."""
    headers = _graph_headers()
    url = (
        f"{GRAPH_ROOT}/me/mailFolders/inbox/messages?"
        "$select=id,subject,sender,toRecipients,body,conversationId,receivedDateTime&"
        f"$filter=receivedDateTime ge {since_iso}&"
        "$orderby=receivedDateTime desc"
    )
    messages = []
    while url:
        response = requests.get(url, headers=headers, verify=False)
        if response.status_code != 200:
            raise RuntimeError(f"Graph API Error: {response.status_code} - {response.text}")
        data = response.json()
        messages.extend(data.get("value", []))
        url = data.get("@odata.nextLink")
    return messages


def fetch_mail_volume(folder, timestamp_field, start_iso, end_iso, page_cb=None):
    """Fetch just enough of a mail folder to count it: ids, timestamps, subjects
    and the counterparty. No bodies.

    `folder` is a well-known folder id ('inbox' / 'sentitems') and
    `timestamp_field` the one that matters for it ('receivedDateTime' /
    'sentDateTime'). `end_iso` is EXCLUSIVE.

    Bodies are the expensive part of a message, so omitting them is what makes a
    multi-year volume pull practical — $top is raised to 999 for the same reason,
    since counting a year of mail is thousands of records and the default page
    size of 10 would be hundreds of round trips. `page_cb(n_so_far)` is called
    after each page so a long pull can show progress.
    """
    headers = _graph_headers()
    messages = []
    # Same treatment as the mailbox-wide fetch: no $orderby (the caller counts,
    # it does not read in order), a modest page size, bounded windows, and
    # retries — a folder-scoped query is cheaper but still times out on a wide
    # range in a busy mailbox.
    for window_start, window_end in _split_window(start_iso, end_iso):
        url = (
            f"{GRAPH_ROOT}/me/mailFolders/{folder}/messages?"
            f"$select=id,subject,conversationId,from,sender,toRecipients,"
            f"{timestamp_field}&"
            f"$filter={timestamp_field} ge {window_start} and "
            f"{timestamp_field} lt {window_end}&"
            f"$top={VOLUME_PAGE_SIZE}"
        )
        while url:
            data = _graph_get(url, headers).json()
            messages.extend(data.get("value", []))
            if page_cb:
                page_cb(len(messages))
            url = data.get("@odata.nextLink")
    return messages


# --------------------------------------------------------------------------- #
# Graph resilience
#
# A mailbox-wide query is expensive server-side, and Graph answers a request it
# cannot serve in time with 502/504 rather than with a smaller page. These are
# transient and a retry usually succeeds, so the volume fetches go through
# _graph_get rather than a bare requests.get.
#
# 429 is throttling, not failure: Graph says when to come back in Retry-After
# and that must be honoured, or retrying makes the throttle worse.
# --------------------------------------------------------------------------- #
GRAPH_RETRY_STATUS = frozenset({429, 500, 502, 503, 504})
GRAPH_MAX_ATTEMPTS = 5
GRAPH_MAX_BACKOFF = 30  # seconds

# Page size for the volume fetches. Far below the documented ceiling of 1000
# deliberately: the docs warn that large pages risk a gateway timeout, and a
# mailbox-wide filter is exactly the query that trips it. More round trips is a
# cheap price for a fetch that finishes.
VOLUME_PAGE_SIZE = 200

# Longest window sent to Graph in one query. A filter spanning years makes the
# server scan the whole mailbox at once; bounded windows keep each request
# small, and let a wide pull report progress instead of appearing to hang.
VOLUME_WINDOW_DAYS = 30


def _graph_get(url, headers=None, attempts=GRAPH_MAX_ATTEMPTS, allow_404=False):
    """GET a Graph URL, retrying throttling and transient gateway errors.

    Raises RuntimeError with the last response body once the attempts are used
    up, or immediately for an error that retrying cannot fix (401, 403, 404).
    With `allow_404`, a 404 returns None instead — some callers ask about things
    that may legitimately not exist.
    """
    headers = headers or _graph_headers()
    backoff = 2.0
    response = None
    for attempt in range(1, attempts + 1):
        response = requests.get(url, headers=headers, verify=False)
        if response.status_code == 200:
            return response
        if response.status_code == 404 and allow_404:
            return None
        if response.status_code not in GRAPH_RETRY_STATUS or attempt == attempts:
            break
        retry_after = response.headers.get("Retry-After", "")
        try:
            wait = float(retry_after)
        except ValueError:
            wait = backoff
        time.sleep(min(wait, GRAPH_MAX_BACKOFF))
        backoff = min(backoff * 2, GRAPH_MAX_BACKOFF)

    status = response.status_code if response is not None else "no response"
    body = (response.text or "")[:400] if response is not None else ""
    hint = ""
    if status in GRAPH_RETRY_STATUS:
        hint = (
            f" Graph kept returning {status} after {attempts} attempts — it is "
            f"struggling with the size of this request. Try a narrower date "
            f"range."
        )
    raise RuntimeError(f"Graph API Error: {status} - {body}{hint}")


def _split_window(start_iso, end_iso, max_days=VOLUME_WINDOW_DAYS):
    """Break an ISO window into consecutive sub-windows of at most `max_days`.

    Yields (start, end) ISO pairs; `end` stays exclusive throughout, so the
    sub-windows tile the range without overlapping or dropping a boundary.
    """
    start = datetime.strptime(start_iso[:10], "%Y-%m-%d")
    end = datetime.strptime(end_iso[:10], "%Y-%m-%d")
    if end <= start:
        yield start_iso, end_iso
        return
    cursor = start
    while cursor < end:
        chunk_end = min(cursor + timedelta(days=max_days), end)
        yield (
            f"{cursor.strftime('%Y-%m-%d')}T00:00:00Z",
            f"{chunk_end.strftime('%Y-%m-%d')}T00:00:00Z",
        )
        cursor = chunk_end


def fetch_well_known_folder_id(well_known_name):
    """The id of a well-known mail folder ('sentitems', 'junkemail', …), or None
    if this mailbox has no such folder.

    Addressed by well-known NAME rather than by matching displayName: display
    names are localized and user-renameable, so name-matching breaks on a
    non-English mailbox. 404 is a legitimate answer — 'clutter' and the hidden
    failure folders do not exist in every mailbox.
    """
    response = _graph_get(
        f"{GRAPH_ROOT}/me/mailFolders/{well_known_name}?$select=id",
        allow_404=True,
    )
    return response.json().get("id") if response is not None else None


def fetch_mailbox_messages(timestamp_field, start_iso, end_iso, page_cb=None):
    """Fetch every message in the MAILBOX for a window — not just one folder.

    This is what makes the volume count honest for a mailbox driven by rules:
    /me/messages spans the whole mailbox (the Graph docs say "including the
    Deleted Items and Clutter folders"), so mail a cloud-side rule filed into a
    custom or nested folder is still counted. Scoping to /mailFolders/inbox
    misses all of it.

    `parentFolderId` comes back with each message so the caller can classify it
    (see email_volume.classify). Enumerating folders and summing them would be
    the obvious alternative and is worse: the folder list includes virtual mail
    SEARCH folders, so any message a search folder matches would be counted
    twice. A mailbox-wide query returns each message once.

    `end_iso` is EXCLUSIVE. Bodies are never selected, which is what keeps a
    multi-year pull practical.
    """
    headers = _graph_headers()
    messages = []
    # No $orderby. Sorting a mailbox-wide filtered result is a large chunk of
    # the server-side work that makes this query time out, and the caller is
    # COUNTING — it buckets by timestamp itself, so the order Graph returns rows
    # in is irrelevant.
    for window_start, window_end in _split_window(start_iso, end_iso):
        url = (
            f"{GRAPH_ROOT}/me/messages?"
            f"$select=id,subject,conversationId,from,sender,toRecipients,"
            f"parentFolderId,{timestamp_field}&"
            f"$filter={timestamp_field} ge {window_start} and "
            f"{timestamp_field} lt {window_end}&"
            f"$top={VOLUME_PAGE_SIZE}"
        )
        while url:
            data = _graph_get(url, headers).json()
            messages.extend(data.get("value", []))
            if page_cb:
                page_cb(len(messages))
            url = data.get("@odata.nextLink")
    return messages


def search_emails_by_subject(subject, start_date, end_date, max_messages=15):
    """Search Outlook for messages whose subject matches a free-text topic within
    the date range. Used by the Executive Translator to scrub related email
    context for the subject being summarized. Returns a list of message dicts."""
    headers = _graph_headers()
    start_str = start_date.strftime("%Y-%m-%d")
    end_str = end_date.strftime("%Y-%m-%d")
    # Escape the WHOLE assembled KQL (incl. the quotes wrapping the subject value)
    # before embedding it in $search="...", matching fetch_emails_for_person.
    kql = f'received:{start_str}..{end_str} AND subject:"{subject}"'.replace('"', '\\"')
    url = (
        f"{GRAPH_ROOT}/me/messages?"
        f'$search="{kql}"&'
        "$select=id,subject,sender,body,receivedDateTime&"
        f"$top={max_messages}"
    )
    response = requests.get(url, headers=headers, verify=False)
    if response.status_code != 200:
        raise RuntimeError(f"Graph API Error: {response.status_code} - {response.text}")
    return response.json().get("value", [])


def fetch_emails_for_person(person_email, start_date, end_date, max_messages=25):
    """Search the mailbox for messages involving a person (sender or recipient)
    within the date range. Used by the 1:1 Meeting Prep Assistant. Returns
    cleaned dicts: {id, subject, sender_name, received, body}."""
    if not person_email:
        return []
    headers = _graph_headers()
    start_str = start_date.strftime("%Y-%m-%d")
    end_str = end_date.strftime("%Y-%m-%d")
    safe = person_email.replace('"', '\\"')
    kql = f'received:{start_str}..{end_str} AND participants:"{safe}"'.replace('"', '\\"')
    url = (
        f"{GRAPH_ROOT}/me/messages?"
        f'$search="{kql}"&'
        "$select=id,subject,sender,receivedDateTime,body&"
        f"$top={max_messages}"
    )
    response = requests.get(url, headers=headers, verify=False)
    if response.status_code != 200:
        raise RuntimeError(f"Graph API Error: {response.status_code} - {response.text}")
    out = []
    for msg in response.json().get("value", []):
        email = (msg.get("sender", {}) or {}).get("emailAddress", {}) or {}
        out.append(
            {
                "id": msg.get("id"),
                "subject": msg.get("subject", "No Subject"),
                "sender_name": email.get("name", "Unknown"),
                "received": msg.get("receivedDateTime", ""),
                "body": clean_html(msg.get("body", {}).get("content", "")),
            }
        )
    return out


def create_mail_draft(to, subject, body_text):
    """Create an Outlook *draft* (does NOT send). Returns the Graph message id.
    This is the HITL-safe default for delegation/follow-up emails."""
    headers = _graph_headers()
    recipients = [to] if isinstance(to, str) else list(to)
    payload = {
        "subject": subject,
        "body": {"contentType": "Text", "content": body_text},
        "toRecipients": [
            {"emailAddress": {"address": r}} for r in recipients if r
        ],
    }
    response = requests.post(
        f"{GRAPH_ROOT}/me/messages", headers=headers, json=payload, verify=False
    )
    response.raise_for_status()
    return response.json().get("id")


def send_mail(to, subject, body_text, draft=True):
    """Send mail via Graph. With draft=True (default) only a draft is created
    so nothing leaves the mailbox without an explicit user action. Set
    draft=False only from a confirmed 'Send' click."""
    if draft:
        return create_mail_draft(to, subject, body_text)

    headers = _graph_headers()
    recipients = [to] if isinstance(to, str) else list(to)
    payload = {
        "message": {
            "subject": subject,
            "body": {"contentType": "Text", "content": body_text},
            "toRecipients": [
                {"emailAddress": {"address": r}} for r in recipients if r
            ],
        },
        "saveToSentItems": True,
    }
    response = requests.post(
        f"{GRAPH_ROOT}/me/sendMail", headers=headers, json=payload, verify=False
    )
    response.raise_for_status()
    return "sent"


# --------------------------------------------------------------------------- #
# MS Graph — Teams & Calendar
# --------------------------------------------------------------------------- #
def post_teams_message(chat_id, text):
    """Post a message to a Teams 1:1/group chat. HITL-gated: call only from a
    confirmed 'Send' click."""
    headers = _graph_headers()
    payload = {"body": {"contentType": "text", "content": text}}
    response = requests.post(
        f"{GRAPH_ROOT}/chats/{chat_id}/messages",
        headers=headers,
        json=payload,
        verify=False,
    )
    response.raise_for_status()
    return response.json().get("id")


def get_calendar_events(start, end):
    """Return calendar events between two datetimes (OoO / calendar features).
    `start`/`end` are datetime objects."""
    headers = _graph_headers()
    params = {
        "startDateTime": start.isoformat(),
        "endDateTime": end.isoformat(),
        "$select": "subject,organizer,start,end,isAllDay,showAs",
        "$orderby": "start/dateTime",
        "$top": "100",
    }
    response = requests.get(
        f"{GRAPH_ROOT}/me/calendarView",
        headers=headers,
        params=params,
        verify=False,
    )
    if response.status_code != 200:
        raise RuntimeError(
            f"Graph calendar error: {response.status_code} - {response.text}"
        )
    return response.json().get("value", [])


# --------------------------------------------------------------------------- #
# Confluence (REST, Bearer PAT)
# --------------------------------------------------------------------------- #
def get_confluence_context(url):
    """Fetch and clean a Confluence page body given its URL (used for the
    StaffMeetingBuilder dedup against last month's notes)."""
    if not url:
        return "None provided."
    match = re.search(r"pageId=(\d+)", url) or re.search(r"pages/(\d+)", url)
    if not match:
        return "None provided."
    page_id = match.group(1)
    token = config.read_token_file(
        config.load_config().get("confluence_pat_path", "")
    )
    api_url = (
        f"{config.CONFLUENCE_DOMAIN}/rest/api/content/{page_id}?expand=body.storage"
    )
    headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
    try:
        response = requests.get(api_url, headers=headers, verify=False)
        response.raise_for_status()
        raw = response.json().get("body", {}).get("storage", {}).get("value", "")
        return clean_html(raw)
    except Exception:
        return "None provided."


# --------------------------------------------------------------------------- #
# Jira (REST v2, Bearer PAT) — extends the ECRT-Test-Plan-Gen pattern
# --------------------------------------------------------------------------- #
def _jira_pat():
    return config.read_token_file(config.load_config().get("jira_pat_path", ""))


def fetch_jira_ticket(ticket_id):
    """Fetch a single Jira issue's summary/description/comments as a text
    block. Mirrors the reference jira_client.fetch_jira_ticket."""
    pat = _jira_pat()
    if not pat:
        return "[Error: Jira PAT not configured.]"
    headers = {"Authorization": f"Bearer {pat}", "Accept": "application/json"}
    api_endpoint = f"{config.JIRA_BASE_URL.rstrip('/')}/rest/api/2/issue/{ticket_id}"
    try:
        response = requests.get(api_endpoint, headers=headers, timeout=15, verify=False)
        if response.status_code == 200:
            fields = response.json().get("fields", {})
            summary = fields.get("summary", "No Summary")
            description = fields.get("description", "No Description")
            comments = fields.get("comment", {}).get("comments", [])
            comments_text = "".join(
                f"\n- {c.get('author', {}).get('displayName', 'Unknown')}: {c.get('body', '')}"
                for c in comments
            )
            return (
                f"JIRA TICKET: {ticket_id}\n\n"
                f"SUMMARY: {summary}\n\n"
                f"DESCRIPTION: {description}\n\n"
                f"COMMENTS: {comments_text}"
            )
        if response.status_code in (401, 403):
            return "[Error: Jira authentication failed. Check your PAT.]"
        if response.status_code == 404:
            return f"[Error: Jira ticket {ticket_id} not found.]"
        return f"[Error: Jira API status {response.status_code}. {response.text}]"
    except Exception as e:
        return f"[Error connecting to Jira: {e}]"


def _stringify_jira_field(value):
    """Best-effort flatten a Jira field value to a short string. Returns "" for
    empty/complex values we can't usefully render."""
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, bool):
        return "Yes" if value else "No"
    if isinstance(value, (int, float)):
        return str(value)
    if isinstance(value, dict):
        for k in ("value", "name", "displayName"):
            v = value.get(k)
            if isinstance(v, str) and v.strip():
                return v.strip()
        return ""
    if isinstance(value, list):
        parts = [_stringify_jira_field(v) for v in value]
        return ", ".join(p for p in parts if p)
    return ""


# Base fields captured explicitly elsewhere, plus noisy collections to skip when
# gathering "other" fields (which surface custom fields like Root Cause Description).
_JIRA_SKIP_FIELDS = {
    "summary", "description", "status", "assignee", "created", "updated",
    "comment", "attachment", "worklog", "subtasks", "issuelinks", "watches",
    "thumbnail", "project", "reporter", "creator", "votes", "progress",
    "aggregateprogress", "timetracking",
}


def fetch_jira_ticket_full(ticket_id):
    """Fetch a single Jira issue as a structured dict for the Phase 3 state
    tracker: metadata, the full comment history with timestamps, attachment
    metadata, and all other non-empty fields (incl. custom fields such as
    'Root Cause Description', labelled via the names map). Raises RuntimeError
    on auth/not-found/API errors so the UI can surface them."""
    pat = _jira_pat()
    if not pat:
        raise RuntimeError("Jira PAT not configured. Add Jira_PAT.txt to the project root.")
    headers = {"Authorization": f"Bearer {pat}", "Accept": "application/json"}
    api_endpoint = f"{config.JIRA_BASE_URL.rstrip('/')}/rest/api/2/issue/{ticket_id}"
    # expand=names returns a {fieldId: displayName} map so custom fields are labelled.
    response = requests.get(
        api_endpoint, headers=headers, params={"expand": "names"}, timeout=30, verify=False
    )
    if response.status_code in (401, 403):
        raise RuntimeError("Jira authentication failed. Check your PAT.")
    if response.status_code == 404:
        raise RuntimeError(f"Jira ticket {ticket_id} not found.")
    if response.status_code != 200:
        raise RuntimeError(f"Jira API status {response.status_code}: {response.text}")

    data = response.json()
    fields = data.get("fields", {})
    names = data.get("names", {})
    extra_fields = []
    for fid, value in fields.items():
        if fid in _JIRA_SKIP_FIELDS:
            continue
        text = _stringify_jira_field(value)
        if text:
            extra_fields.append({"name": names.get(fid, fid), "value": text})
    comments = [
        {
            "id": c.get("id", ""),
            "author": (c.get("author") or {}).get("displayName", "Unknown"),
            "created": c.get("created", ""),
            "body": c.get("body", ""),
        }
        for c in fields.get("comment", {}).get("comments", [])
    ]
    attachments = [
        {
            "id": a.get("id"),
            "filename": a.get("filename", ""),
            "mime": a.get("mimeType", ""),
            "size": a.get("size", 0),
            "content_url": a.get("content", ""),
            "created": a.get("created", ""),
        }
        for a in fields.get("attachment", []) or []
    ]
    return {
        "key": data.get("key", ticket_id),
        "summary": fields.get("summary", ""),
        "description": fields.get("description", "") or "",
        "status": (fields.get("status") or {}).get("name", ""),
        "assignee": (fields.get("assignee") or {}).get("displayName", "Unassigned"),
        "created": fields.get("created", ""),
        "updated": fields.get("updated", ""),
        "comments": comments,
        "attachments": attachments,
        "extra_fields": extra_fields,
    }


def fetch_carrier_tracking_text(carrier, tracking_number):
    """Best-effort fetch of a carrier's public tracking page; returns
    (cleaned_text, url). Many carrier pages are JavaScript shells, so the text may
    be empty/uninformative — the caller falls back to 'Unknown' + the link."""
    url = config.carrier_tracking_url(carrier, tracking_number)
    if not url:
        return "", ""
    try:
        response = requests.get(
            url, headers={"User-Agent": "Mozilla/5.0"}, verify=False, timeout=15
        )
        if not response.ok:
            return "", url
        return clean_html(response.text)[:6000], url
    except Exception:
        return "", url


def download_jira_attachment(content_url):
    """Download a Jira attachment's bytes from its content URL (Bearer PAT)."""
    pat = _jira_pat()
    if not pat:
        raise RuntimeError("Jira PAT not configured.")
    headers = {"Authorization": f"Bearer {pat}"}
    response = requests.get(content_url, headers=headers, timeout=30, verify=False)
    response.raise_for_status()
    return response.content


def search_jira(jql, max_results=50, fields=None):
    """Run a JQL search and return a simplified list of issue dicts:
    {key, summary, status, assignee, updated}. Raises RuntimeError with a
    readable message on failure so the UI can surface it cleanly."""
    pat = _jira_pat()
    if not pat:
        raise RuntimeError("Jira PAT not configured. Add Jira_PAT.txt to the project root.")
    if fields is None:
        fields = ["summary", "status", "assignee", "updated"]
    headers = {"Authorization": f"Bearer {pat}", "Accept": "application/json"}
    api_endpoint = f"{config.JIRA_BASE_URL.rstrip('/')}/rest/api/2/search"
    params = {"jql": jql, "maxResults": max_results, "fields": ",".join(fields)}
    response = requests.get(
        api_endpoint, headers=headers, params=params, timeout=20, verify=False
    )
    if response.status_code in (401, 403):
        raise RuntimeError("Jira authentication failed. Check your PAT.")
    if response.status_code != 200:
        raise RuntimeError(f"Jira API status {response.status_code}: {response.text}")

    issues = []
    for issue in response.json().get("issues", []):
        f = issue.get("fields", {})
        assignee = f.get("assignee")
        issues.append(
            {
                "key": issue.get("key"),
                "summary": f.get("summary", ""),
                "status": (f.get("status") or {}).get("name", ""),
                "assignee": (assignee or {}).get("displayName", "Unassigned"),
                "updated": (f.get("updated") or "")[:10],
            }
        )
    return issues
